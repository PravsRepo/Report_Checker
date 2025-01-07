import pdfplumber
import pandas as pd
import os

from difflib import get_close_matches
from pptx import Presentation
from docx import Document


class ReportChecker:

    def __init__(self, file_name, input_path):
        self.file_name = file_name
        self.input_path = input_path

    # read excel sheet to get students name and their team number
    def read_file(self):
        df = pd.read_excel(
            self.file_name, sheet_name="Sheet1")
        return df

    def get_report_extn(self, df):
        starting_row = 1
        file_extn = os.listdir(self.input_path)
        for file in file_extn:
            raw_name, extension = os.path.splitext(file)
            team_num = raw_name.replace("Team ", "")
            if extension == ".pptx":
                to_write_data = self.read_ppt(f"{input_path}/{file}", team_num)
            elif extension == ".pdf":
                to_write_data = self.read_pdf(f"{input_path}/{file}", team_num)
            elif extension == "docx":
                self.read_word(f"{input_path}/{file}", team_num)
            else:
                print(f"file name: {file}", team_num)

    def read_ppt(self, source_file, team_num):
        prs = Presentation(source_file)
        text_chunks = []
        text_runs = {}
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_chunks.append(cell.text)
            else:
                None
        text_runs[team_num] = text_chunks
        df = self.read_file()
        output = self.check_names(df, text_runs)
        return output

    def read_pdf(self, source_file, team_num):
        text_runs_pdf = {}
        with pdfplumber.open(source_file) as pdf:
            first_page = pdf.pages[0]
            text_chunks_pdf = first_page.extract_text()
        text_runs_pdf[team_num] = text_chunks_pdf
        df = self.read_file()
        output = self.verify_names_with_run(df, text_runs_pdf)
        return output

    def read_word(self):
        pass

    # check student names in the report output and return the output - PPT
    def check_names(self, df, text_runs):
        team_num = next(iter(text_runs))
        team_df = df[df["Teams"] == int(team_num)]
        name_df = team_df["Name"]

        stop_words = {"Student Name", "Year of Study", "Department", "ECE", "CSE",
                      "IT", "S. No.", "Roll #", "AIDS", "MECHANICAL", "Faculty Name"}
        text_runs_values = text_runs.get(team_num)
        result = []
        for run in text_runs_values:
            clean_run = run.strip()
            clean_run_lower = clean_run.lower()
            stop_words_lower = {word.lower() for word in stop_words}
            if (clean_run_lower in stop_words_lower or clean_run.isnumeric() or clean_run == ''
                    or len(clean_run) <= 3 or any(char in clean_run for char in [":", "-", "/", "&", ")", "("])):
                None
            else:
                result.append(clean_run)
        print(f"{team_num}: {result}")

        # First set of results - cross check the names with the dataframe
        pattern = "|".join(result)
        first_set_results = team_df["Name"].str.contains(pattern, case=False, na=False)
        first_set_results.name = "First set of results"
        first_set_results.reset_index(drop=True, inplace=True)

        # second set of results
        pattern = "|". join(result).lower()
        matches = []
        for name in team_df["Name"]:
            iter_name = name.lower().split()
            found = False
            for i in iter_name:
                if len(i)>1:
                    if i in pattern:
                        found = True
                        break
            matches.append(found)
        second_set_results = pd.DataFrame(matches, columns=["Second set of results"])

        # Third set of results 
        close_results = []
        for name in name_df:
            close_results.append(bool(get_close_matches(name, list(result))))
        third_set_results = pd.DataFrame(close_results, columns=["Third set of results"])

        # finalize the results
        second_set_results = second_set_results.squeeze()
        third_set_results = third_set_results.squeeze()

        final_results = first_set_results | second_set_results | third_set_results
        return final_results


    # check student names in the report output and return the output - PDF
    def verify_names_with_run(self, df, text_runs_pdf):
        team_num = next(iter(text_runs_pdf))
        team_df = df[df["Teams"] == int(team_num)]
        name_df = team_df["Name"]

        result = []
        stop_words = {"Student", "Year", "year", "of", "Study", "Department", "IT", "S.No.", "Date", "LPB01", "Presentation", "B.Tech", "Team", "College", "LEAP.",
                      "Roll", "AIDS", "MECHANICAL", "EEE", "CSE", "Faculty", "Name", "#", "Engineering", "Reverse", "Report", "Product", "Copyright", "reserved", "rights"}
        text_runs_values = text_runs_pdf.get(team_num).split()
        for run in text_runs_values:
            clean_run = run.strip()
            clean_run_lower = clean_run.lower()
            stop_words_lower = {word.lower() for word in stop_words}
            if (clean_run_lower in stop_words_lower or clean_run.isnumeric() or clean_run == ''
                    or len(clean_run) <= 3 or any(char in clean_run for char in [":", "-", "/", "&", ")", "("])):
                None
            else:
                result.append(clean_run)
        print(f"{team_num}: {result}")

        # first set of results
        pattern = "|".join(result)
        first_set_results = team_df["Name"].str.contains(pattern, case=False, na=False)
        first_set_results.name = "First set of results"
        first_set_results.reset_index(drop=True, inplace=True)

        # second set of results
        pattern = "|". join(result).lower()
        matches = []
        for name in team_df["Name"]:
            iter_name = name.lower().split()
            found = False
            for i in iter_name:
                if len(i)>1:
                    if i in pattern:
                        found = True
                        break
            matches.append(found)
        second_set_results = pd.DataFrame(matches, columns=["Second set of results"])


        # third set of results - get close matches
        close_results = []
        for name in name_df:
            close_results.append(bool(get_close_matches(name, list(result))))
        third_set_results = pd.DataFrame(close_results, columns=["Third set of results"])

        # finalize the results
        second_set_results = second_set_results.squeeze()
        third_set_results = third_set_results.squeeze()

        final_results = first_set_results| second_set_results | third_set_results
        return final_results

    def write_excel(self, to_write_data, starting_row, team_num):
        file_path = "D:/Report_checker/output/LPB01-B02-Prince BC Evaluation V0.6_output.xlsx"
        with pd.ExcelWriter(file_path, engine="openpyxl", mode="a", if_sheet_exists="overlay") as writer:
            to_write_data.to_excel(excel_writer=writer,
                                   sheet_name="Sheet1", startcol=3, startrow=starting_row, index=False, header=False)
        print(f"{team_num}: Data written successful...\n\n")


file_name = "D:/Report_checker/input/LPB01-B02-Prince BC Evaluation V0.6.xlsx"
input_path = "D:/Report_checker/input/Testing_folder"
checker_obj = ReportChecker(file_name, input_path)
df = checker_obj.read_file()
to_write_data = checker_obj.get_report_extn(df)

