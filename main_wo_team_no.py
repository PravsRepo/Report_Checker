import pdfplumber
import pandas as pd
import os

from difflib import get_close_matches
from pptx import Presentation



class ReportFinder:

    def __init__(self, file_name, input_path):
        self.file_name = file_name
        self.input_path = input_path

    # read excel sheet to get students name and their team number
    def read_file(self):
        df = pd.read_excel(self.file_name, sheet_name="Consolidated")
        return df

    def get_report_extn(self):
        file_extn = os.listdir(self.input_path)
        for file in file_extn:
            raw_name, extension = os.path.splitext(file)
            if extension == ".pptx":
               self.read_ppt(f"{input_path}/{file}", raw_name)
            elif extension == ".pdf":
                self.read_pdf(f"{input_path}/{file}", raw_name)
            elif extension == "docx":
                self.read_word(f"{input_path}/{file}")
            else:
                print(f"file name: {file}")

    def read_ppt(self, source_file, raw_name):
        prs = Presentation(source_file)
        text_runs = []
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_runs.append(cell.text)
            else:
                None
        self.clean_text_runs(text_runs, raw_name)
        

    def read_pdf(self, source_file, raw_name):
        text_runs_pdf = []
        with pdfplumber.open(source_file) as pdf:
            first_page = pdf.pages[0]
            result = first_page.extract_table()
        for i in result:
            text_runs_pdf.extend(i)
        self.clean_text_runs(text_runs_pdf, raw_name)
        
    
    def clean_text_runs(self, text_runs, raw_name):
        stop_words = {"Student Name", "Year of Study", "Department", "S. No.", "Roll #", "Faculty Name", "First year", "1st Year"}
        result = []
        for run in text_runs:
            if run is not None:
                clean_run = run.strip()
                clean_run_lower = clean_run.lower()
                stop_words_lower = {word.lower() for word in stop_words}
                if (clean_run_lower in stop_words_lower or clean_run.isalnum() or clean_run.isnumeric() or clean_run == ''
                        or len(clean_run) <= 3 or any(char in clean_run for char in [":", "-", "/", "&", ")", "("])):
                    None
                else:
                    result.append(clean_run)
        self.check_names(result)

    def check_names(self, result):
        df = self.read_file()
        name_df = df["Student Name"]

        # First set of results - cross check the names with the dataframe
        pattern = "|".join(result)
        first_set_results = df["Student Name"].str.contains(pattern, case=False, na=False)
        first_set_results.name = "First set of results"
        first_set_results = first_set_results.loc[first_set_results]

        # second set of results
        pattern = "|". join(result).lower()
        matches = []
        for name in df["Student Name"]:
            iter_name = name.lower().split()
            found = False
            for i in iter_name:
                if len(i)>1:
                    if i in pattern:
                        found = True
                        break
            matches.append(found)
        second_set_results = pd.Series(matches, name="Second set of results")
        second_set_results = second_set_results.loc[second_set_results]

        # Third set of results 
        close_results = []
        for name in name_df:
            close_results.append(bool(get_close_matches(name, list(result))))
        third_set_results = pd.Series(close_results, name="Third set of results")
        third_set_results = third_set_results.loc[third_set_results]

        final_results = first_set_results | second_set_results | third_set_results
        print(f"Final results: {final_results}")


    def write_excel(self, to_write_data):
        file_path = "D:/Report_checker/output/LPB01-B02-Prince BC Evaluation V0.6_output.xlsx"
        with pd.ExcelWriter(file_path, engine="openpyxl", mode="a", if_sheet_exists="overlay") as writer:
            to_write_data.to_excel(excel_writer=writer,
                                   sheet_name="Sheet1", startcol=2, startrow=1, index=False, header=False)
        print(f"Data written successful...\n\n")
   

file_name = "D:/Report_checker/input/LPB02-ICT student BC eval v1.0.xlsx"
input_path = "D:/Report_checker/input/Test files"
finder_obj = ReportFinder(file_name, input_path)
finder_obj.read_file()
finder_obj.get_report_extn()