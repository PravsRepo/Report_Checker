from pptx import Presentation

prs = Presentation("D:/Report_checker/input/team 51.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    # print(slide)
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for cell in table.iter_cells():
                text_runs.append(cell.text)
    break
print(text_runs)